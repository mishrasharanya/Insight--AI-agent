import os
import json
import hashlib
from datetime import datetime
from pathlib import Path

import pandas as pd
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from supabase_client import get_client

COLLECTION_NAME = "personal_memory"
DATA_FOLDER = "data"
DATA_SOURCES_FILE = "data_sources.json"
MAX_CHUNK_CHARS = 800

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".csv", ".pdf",
    ".docx", ".pptx", ".xlsx", ".json"
}

# File types that also get structured-row extraction for charting,
# in addition to the normal text chunking every file type gets.
STRUCTURED_EXTENSIONS = {".csv", ".xlsx"}

embedding_model = SentenceTransformer("all-MiniLM-L6-v2")


def load_data_folders():
    if not Path(DATA_SOURCES_FILE).exists():
        return [Path(DATA_FOLDER)]

    with open(DATA_SOURCES_FILE, "r", encoding="utf-8") as file:
        config = json.load(file)

    folders = config.get("folders", [DATA_FOLDER])
    return [Path(folder).expanduser() for folder in folders]


def get_file_hash(file_path):
    with open(file_path, "rb") as file:
        return hashlib.md5(file.read()).hexdigest()


def read_text_file(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()


def read_csv_file(file_path):
    df = pd.read_csv(file_path)
    rows = []

    for _, row in df.iterrows():
        row_text = [f"{column}: {row[column]}" for column in df.columns]
        rows.append(" | ".join(row_text))

    return "\n".join(rows)


def read_pdf_file(file_path):
    reader = PdfReader(file_path)
    pages = []

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {page_num}]\n{text.strip()}")

    return "\n\n".join(pages)


def read_docx_file(file_path):
    document = Document(file_path)
    parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def read_pptx_file(file_path):
    presentation = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(presentation.slides, start=1):
        parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)

        if parts:
            slides.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

    return "\n\n".join(slides)


def read_xlsx_file(file_path):
    workbook = load_workbook(file_path, data_only=True)
    sheets = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows = []

        for row in sheet.iter_rows(values_only=True):
            values = [
                str(value).strip()
                for value in row
                if value is not None and str(value).strip()
            ]

            if values:
                rows.append(" | ".join(values))

        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    return "\n\n".join(sheets)


def read_json_file(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        data = json.load(file)

    return json.dumps(data, indent=2, ensure_ascii=False)


def read_file(file_path):
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    if suffix in [".txt", ".md"]:
        return read_text_file(file_path)

    if suffix == ".csv":
        return read_csv_file(file_path)

    if suffix == ".pdf":
        return read_pdf_file(file_path)

    if suffix == ".docx":
        return read_docx_file(file_path)

    if suffix == ".pptx":
        return read_pptx_file(file_path)

    if suffix == ".xlsx":
        return read_xlsx_file(file_path)

    if suffix == ".json":
        return read_json_file(file_path)

    return ""


def load_document(file_path):
    return read_file(Path(file_path))


# ---------- Structured row extraction (new - for charts) ----------
# NOTE: kept deliberately separate from read_csv_file/read_xlsx_file above.
# Those two produce flattened text for embedding/chunking - this produces
# typed row data for the structured_rows table, which a charts endpoint
# can query directly without going through the LLM or chunked text at all.

def extract_csv_rows(file_path):
    df = pd.read_csv(file_path)
    rows = []

    for index, row in df.iterrows():
        row_dict = {}
        for column in df.columns:
            value = row[column]
            # Keep JSON-serializable types only - pandas/numpy scalars
            # don't serialize cleanly otherwise.
            if pd.isna(value):
                row_dict[column] = None
            elif isinstance(value, (int, float, str, bool)):
                row_dict[column] = value
            else:
                row_dict[column] = str(value)

        rows.append({
            "sheet_name": None,
            "row_index": index,
            "row_data": row_dict,
        })

    return rows


def extract_xlsx_rows(file_path):
    workbook = load_workbook(file_path, data_only=True)
    all_rows = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows_iter = sheet.iter_rows(values_only=True)

        try:
            header = next(rows_iter)
        except StopIteration:
            continue

        header = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(header)]

        for row_index, row in enumerate(rows_iter):
            if all(value is None for value in row):
                continue

            row_dict = {}
            for column, value in zip(header, row):
                if value is None:
                    row_dict[column] = None
                elif isinstance(value, (int, float, str, bool)):
                    row_dict[column] = value
                else:
                    row_dict[column] = str(value)

            all_rows.append({
                "sheet_name": sheet_name,
                "row_index": row_index,
                "row_data": row_dict,
            })

    return all_rows


def extract_structured_rows(file_path):
    suffix = Path(file_path).suffix.lower()

    if suffix == ".csv":
        return extract_csv_rows(file_path)

    if suffix == ".xlsx":
        return extract_xlsx_rows(file_path)

    return []


def chunk_text(text, max_chars=MAX_CHUNK_CHARS):
    paragraphs = text.split("\n")
    chunks = []
    current_chunk = ""

    for paragraph in paragraphs:
        paragraph = paragraph.strip()

        if not paragraph:
            continue

        if len(current_chunk) + len(paragraph) + 1 <= max_chars:
            current_chunk += "\n" + paragraph
        else:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = paragraph

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


def get_effective_date(file_path):
    modified_time = os.path.getmtime(file_path)
    return datetime.fromtimestamp(modified_time).isoformat()


def get_all_files_from_folders(folders):
    files = []

    for folder in folders:
        if not folder.exists():
            print(f"Folder not found, skipping: {folder}")
            continue

        for file_path in folder.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                files.append(file_path)

    return files


def upsert_chunks(client, collection_name, file_path, chunks, file_hash, effective_date):
    """
    Embeds and upserts text chunks into the `chunks` table. Mirrors the
    dedup behavior the old Chroma `collection.upsert()` gave for free -
    the (collection_name, file_hash, chunk_index) unique constraint means
    re-ingesting an unchanged file overwrites rather than duplicates.
    """
    records = []

    for index, chunk in enumerate(chunks):
        embedding = embedding_model.encode(chunk).tolist()

        records.append({
            "collection_name": collection_name,
            "filename": file_path.name,
            "source": str(file_path),
            "file_type": file_path.suffix.lower(),
            "chunk_index": index,
            "file_hash": file_hash,
            "effective_date": effective_date,
            "text": chunk,
            "embedding": embedding,
        })

    if records:
        client.table("chunks").upsert(
            records,
            on_conflict="collection_name,file_hash,chunk_index",
        ).execute()

    return len(records)


def upsert_structured_rows(client, collection_name, file_path, effective_date):
    """
    Extracts and stores typed rows from csv/xlsx files for the charts
    endpoint. Unlike `chunks`, there's no dedup constraint here yet - if
    you re-ingest the same file repeatedly, you'll get duplicate rows.
    Fine for now, but worth revisiting once re-ingestion is a common flow.
    """
    rows = extract_structured_rows(file_path)

    if not rows:
        return 0

    records = [
        {
            "collection_name": collection_name,
            "filename": file_path.name,
            "sheet_name": row["sheet_name"],
            "row_index": row["row_index"],
            "row_data": row["row_data"],
            "effective_date": effective_date,
        }
        for row in rows
    ]

    client.table("structured_rows").insert(records).execute()
    return len(records)


def ingest_files(collection_name=None):
    """
    collection_name: same role it plays everywhere else in the codebase -
    defaults to the shared COLLECTION_NAME for desktop/single-user use,
    or pass a per-user string for Google-synced ingestion (see
    google_sync.collection_name_for).
    """
    collection_name = collection_name or COLLECTION_NAME
    client = get_client()

    folders = load_data_folders()
    files = get_all_files_from_folders(folders)

    print("Reading from folders:")
    for folder in folders:
        print(f"- {folder}")

    print(f"\nFound {len(files)} supported files.\n")

    total_chunks = 0
    total_structured_rows = 0

    for file_path in files:
        print(f"Ingesting: {file_path}")

        try:
            text = read_file(file_path)
        except Exception as error:
            print(f"Failed to read {file_path}: {error}")
            continue

        if not text.strip():
            print(f"No text extracted from: {file_path}")
        else:
            chunks = chunk_text(text)
            file_hash = get_file_hash(file_path)
            effective_date = get_effective_date(file_path)

            count = upsert_chunks(client, collection_name, file_path, chunks, file_hash, effective_date)
            total_chunks += count

        if file_path.suffix.lower() in STRUCTURED_EXTENSIONS:
            try:
                effective_date = get_effective_date(file_path)
                row_count = upsert_structured_rows(client, collection_name, file_path, effective_date)
                total_structured_rows += row_count
            except Exception as error:
                print(f"Failed to extract structured rows from {file_path}: {error}")

    print(f"\nIngestion complete. Total chunks added: {total_chunks}")
    print(f"Total structured rows added: {total_structured_rows}")


if __name__ == "__main__":
    ingest_files()